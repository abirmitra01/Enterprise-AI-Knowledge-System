from pptx import Presentation


def extract_ppt_text(filepath):

    presentation = Presentation(filepath)

    text = ""

    for slide in presentation.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text += shape.text + "\n"

    return text